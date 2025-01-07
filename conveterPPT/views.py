from django.shortcuts import render
from django.http import HttpResponse
from .forms import DocumentUploadForm
from docx import Document
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# The THEMES dictionary
THEMES = {
    "dark": {"bg_color": (30, 30, 30), "title_color": (255, 255, 255), "content_color": (200, 200, 200)},
    "corporate": {"bg_color": (240, 240, 240), "title_color": (0, 0, 100), "content_color": (0, 0, 0)},
    "cyberpunk": {"bg_color": (255, 230, 0), "title_color": (255, 0, 255), "content_color": (0, 255, 255)},
    "cupcake": {"bg_color": (255, 204, 204), "title_color": (255, 51, 153), "content_color": (204, 153, 255)},
    "bumblebee": {"bg_color": (255, 255, 102), "title_color": (255, 204, 0), "content_color": (0, 0, 0)},
    "emerald": {"bg_color": (80, 200, 120), "title_color": (255, 255, 255), "content_color": (0, 0, 0)},
    "synthwave": {"bg_color": (50, 0, 100), "title_color": (255, 0, 255), "content_color": (0, 255, 255)},
    "retro": {"bg_color": (255, 243, 205), "title_color": (255, 102, 0), "content_color": (255, 204, 102)},
    "halloween": {"bg_color": (0, 0, 0), "title_color": (255, 153, 0), "content_color": (255, 0, 0)},
    "garden": {"bg_color": (153, 204, 102), "title_color": (255, 102, 204), "content_color": (102, 51, 51)},
}

def apply_theme(slide, theme_name, title_text, content_text):
    theme = THEMES.get(theme_name, THEMES["dark"])  

    # Set Background Color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*theme["bg_color"])

    # Title Styling
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme["title_color"])

    # Content Styling
    content_shape = slide.shapes.placeholders[1]
    content_shape.text = content_text
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(*theme["content_color"])
        paragraph.alignment = PP_ALIGN.LEFT

def process_document_and_create_ppt(request):
    if request.method == "POST":
        form = DocumentUploadForm(request.POST, request.FILES)
        if form.is_valid():
            uploaded_file = request.FILES['document']
            selected_theme = request.POST.get('theme', 'dark')  

            document = Document(uploaded_file)
            headings = []
            current_heading = None
            content = []

            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text and text.istitle():
                    if current_heading:
                        headings.append({'title': current_heading, 'content': "\n".join(content)})
                    current_heading = text
                    content = []
                elif text:
                    content.append(text)

            if current_heading:
                headings.append({'title': current_heading, 'content': "\n".join(content)})

            presentation = Presentation()
            for heading in headings:
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                apply_theme(slide, selected_theme, heading['title'], heading['content'])

            pptx_file = BytesIO()
            presentation.save(pptx_file)
            pptx_file.seek(0)
            response = HttpResponse(
                pptx_file,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            response["Content-Disposition"] = 'attachment; filename="professional_ppt.pptx"'
            return response
    else:
        form = DocumentUploadForm()

    return render(request, 'docUpload.html', {'form': form, 'themes': THEMES})

